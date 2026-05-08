import os
import asyncio
import requests
import io
import logging
import google.generativeai as genai
from dotenv import load_dotenv
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command
from aiogram.types import FSInputFile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Sozlamalar
load_dotenv()
BOT_TOKEN = os.getenv("BOT_TOKEN")
GEMINI_KEY = os.getenv("GEMINI_API_KEY")

logging.basicConfig(level=logging.INFO)
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()

# --- GEMINI MODELINI TO'G'RI SOZLASH ---
def get_working_model():
    try:
        genai.configure(api_key=GEMINI_KEY)
        models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods]
        # Eng yaxshi modellarni tartib bilan tekshirish
        for m_name in ['models/gemini-1.5-flash', 'models/gemini-1.5-pro', 'models/gemini-pro']:
            if m_name in models:
                return genai.GenerativeModel(m_name)
        return genai.GenerativeModel(models[0]) if models else None
    except:
        return None

model = get_working_model()

# --- TOZA MATN YARATISH ---
async def get_clean_content(topic):
    prompt = (
        f"Mavzu: '{topic}'. Ushbu mavzuda 5 slayidli professional taqdimot rejasi tuz. "
        "Faqat o'zbek tilida javob ber. Har bir slaydni aynan quyidagi formatda yoz, "
        "ortiqcha 'Mana reja' yoki 'Slayd 1' kabi so'zlarni ishlatma:\n"
        "Sarlavha | Matn (qisqa 3 ta qator) | Rasm uchun inglizcha qisqa nom\n"
        "Namuna:\nO'zbekiston turizmi | 1. Qadimiy shaharlar\n2. Milliy taomlar\n3. Mehmondoshlik | Samarkand Registan"
    )
    try:
        if model:
            response = model.generate_content(prompt)
            return response.text
        # Zaxira AI (Pollinations)
        res = requests.get(f"https://text.pollinations.ai/{prompt}?model=llama")
        return res.text
    except:
        return None

# --- SIFATLI RASM OLISH ---
def get_hq_image(keyword):
    # Rasm sifatli chiqishi uchun promptni boyitish
    query = f"professional high quality presentation slide image of {keyword}"
    url = f"https://image.pollinations.ai/prompt/{query.replace(' ', '%20')}?width=1024&height=768&nologo=true"
    try:
        res = requests.get(url, timeout=20)
        if res.status_code == 200:
            return io.BytesIO(res.content)
    except:
        return None

# --- PPTX YARATISH (DIZAYN BILAN) ---
def create_styled_pptx(topic, content, filename):
    prs = Presentation()
    
    # 1. Titul slayd (Chiroyli sarlavha)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = topic.upper()
    subtitle = slide.placeholders[1]
    subtitle.text = "Sun'iy intellekt tomonidan tayyorlangan professional taqdimot"

    lines = [l.strip() for l in content.split('\n') if "|" in l]
    
    for line in lines[:5]:
        try:
            parts = line.split("|")
            s_title = parts[0].strip()
            s_text = parts[1].strip()
            s_img = parts[2].strip() if len(parts) > 2 else s_title

            # Slayd yaratish (Layout 1: Title and Content)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Sarlavha stili
            title_shape = slide.shapes.title
            title_shape.text = s_title
            
            # Matn stili
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = s_text
            
            # RASM JOYLASH (O'ng tomonga, aniq o'lchamda)
            img_stream = get_hq_image(s_img)
            if img_stream:
                # Inches(5.5) - chapdan, Inches(1.2) - tepadan, Inches(4) - eni
                slide.shapes.add_picture(img_stream, Inches(5.2), Inches(1.5), width=Inches(4.5))
        except:
            continue

    prs.save(filename)

# --- BOT INTERFEYSI ---
@dp.message(Command("start"))
async def start(m: types.Message):
    await m.answer(f"Assalomu alaykum {m.from_user.first_name}!\nMen sizga **toza o'zbek tilida** va **rasmlar bilan** professional taqdimot yaratib beraman.\n\nMavzuni yuboring:")

@dp.message(F.text)
async def handle_request(m: types.Message):
    topic = m.text
    status = await m.answer("⏳ **Matn va rasmlar tayyorlanmoqda...**\n(Bu 15-20 soniya vaqt olishi mumkin)")
    
    ai_content = await get_clean_content(topic)
    if not ai_content:
        await status.edit_text("❌ Xatolik: AI matn yarata olmadi.")
        return

    file_name = f"pres_{m.from_user.id}.pptx"
    try:
        create_styled_pptx(topic, ai_content, file_name)
        
        await m.answer_document(
            FSInputFile(file_name), 
            caption=f"✅ **Taqdimot tayyor!**\n\nMavzu: {topic}\nModel: {model.model_name if model else 'Llama'}"
        )
    except Exception as e:
        await m.answer(f"⚠️ Xato: {str(e)}")
    finally:
        if os.path.exists(file_name):
            os.remove(file_name)
        await status.delete()

async def main():
    # Eski so'rovlarni tozalash (ConflictError oldini olish uchun)
    await bot.delete_webhook(drop_pending_updates=True)
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
